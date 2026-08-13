#!/usr/bin/env python3
"""
PPTX Import Script for DeckMason

Reads an existing .pptx file and extracts slide content into a structured
JSON file that DeckMason can use as source material for HTML generation.

Usage:
    python3 pptx_import.py input.pptx [--output decks/input_deck.json]

Extracts per slide:
    - Slide layout type
    - Title text
    - Body text and paragraphs
    - Images with positions and sizes
    - Speaker notes

Extracts from theme1.xml inside the pptx:
    - Theme colors (color scheme)
    - Theme fonts (major and minor font)

Requirements:
    pip3 install python-pptx

Output:
    A JSON file with this structure:
    {
        "source_file": "input.pptx",
        "slide_count": 10,
        "theme": {
            "colors": {"dk1": "#000000", "lt1": "#ffffff", ...},
            "fonts": {"major": "Calibri Light", "minor": "Calibri"}
        },
        "slides": [
            {
                "index": 0,
                "layout": "Title Slide",
                "title": "Presentation Title",
                "subtitle": "Subtitle text",
                "body": ["Bullet 1", "Bullet 2"],
                "paragraphs": [{"text": "...", "level": 0}],
                "images": [{"path": "media/image1.png", "left": 100, "top": 200, "width": 400, "height": 300}],
                "notes": "Speaker notes text",
                "background_color": "#ffffff"
            }
        ]
    }
"""

import argparse
import json
import os
import sys
import tempfile
import zipfile
import xml.etree.ElementTree as ET

# python-pptx is required
try:
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx is not installed.", file=sys.stderr)
    print("Install it with: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)


# XML namespaces used in OOXML theme files
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def emu_to_px(emu_value):
    """Convert EMU (English Metric Units) to pixels at 96 DPI."""
    if emu_value is None:
        return 0
    # 1 inch = 914400 EMU, 1 inch = 96 px at 96 DPI
    return round(emu_value / 914400 * 96)


def rgb_to_hex(rgb):
    """Convert an RGBColor to a hex string like #ff0000."""
    if rgb is None:
        return None
    return "#{:02x}{:02x}{:02x}".format(rgb[0], rgb[1], rgb[2])


def extract_theme_from_pptx(pptx_path):
    """
    Extract theme colors and fonts from theme1.xml inside the pptx archive.

    The theme file lives at ppt/theme/theme1.xml and contains:
    - <a:clrScheme> with dk1, lt1, dk2, lt2, accent1-6
    - <a:fontScheme> with majorFont and minorFont
    """
    theme = {"colors": {}, "fonts": {}}

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            # Find the theme file
            theme_path = None
            for name in zf.namelist():
                if name.startswith("ppt/theme/") and name.endswith(".xml"):
                    theme_path = name
                    break

            if theme_path is None:
                print("WARNING: No theme file found in pptx", file=sys.stderr)
                return theme

            theme_data = zf.read(theme_path)
            root = ET.fromstring(theme_data)

            # Extract color scheme
            clr_scheme = root.find(".//a:clrScheme", NS)
            if clr_scheme is not None:
                color_names = ["dk1", "lt1", "dk2", "lt2",
                               "accent1", "accent2", "accent3", "accent4",
                               "accent5", "accent6", "hlink", "folHlink"]
                for color_name in color_names:
                    elem = clr_scheme.find("a:{}".format(color_name), NS)
                    if elem is not None:
                        # Color can be <a:srgbClr val="RRGGBB"/> or <a:sysClr val="windowText" lastClr="000000"/>
                        srgb = elem.find("a:srgbClr", NS)
                        sysclr = elem.find("a:sysClr", NS)
                        if srgb is not None:
                            hex_val = srgb.get("val", "")
                            if hex_val:
                                theme["colors"][color_name] = "#{}".format(hex_val.lower())
                        elif sysclr is not None:
                            hex_val = sysclr.get("lastClr", "")
                            if hex_val:
                                theme["colors"][color_name] = "#{}".format(hex_val.lower())

            # Extract font scheme
            font_scheme = root.find(".//a:fontScheme", NS)
            if font_scheme is not None:
                major_font = font_scheme.find("a:majorFont/a:latin", NS)
                minor_font = font_scheme.find("a:minorFont/a:latin", NS)
                if major_font is not None:
                    theme["fonts"]["major"] = major_font.get("typeface", "")
                if minor_font is not None:
                    theme["fonts"]["minor"] = minor_font.get("typeface", "")

    except zipfile.BadZipFile:
        print("WARNING: Could not read pptx as zip archive for theme extraction", file=sys.stderr)
    except ET.ParseError:
        print("WARNING: Could not parse theme1.xml", file=sys.stderr)

    return theme


def extract_slide_data(slide, index):
    """Extract all relevant data from a single slide."""
    slide_data = {
        "index": index,
        "layout": "",
        "title": "",
        "subtitle": "",
        "body": [],
        "paragraphs": [],
        "images": [],
        "notes": "",
        "background_color": None,
    }

    # Layout name
    try:
        slide_data["layout"] = slide.slide_layout.name
    except (AttributeError, KeyError):
        slide_data["layout"] = "Unknown"

    # Background color
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None and str(fill.type) != "MSO_FILL.NONE (0)":
            try:
                fore = fill.fore_color
                if fore is not None:
                    slide_data["background_color"] = rgb_to_hex(fore.rgb)
            except (TypeError, AttributeError):
                pass
    except (AttributeError, KeyError, TypeError):
        pass

    # Extract text from all shapes
    for shape in slide.shapes:
        shape_type = str(shape.shape_type) if shape.shape_type else "UNKNOWN"

        # Check if this is the title placeholder
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            if not text:
                continue

            # Identify title vs body by placeholder type or position
            is_title = False
            is_subtitle = False
            try:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    type_str = str(ph_type).lower()
                    # idx 0 = title, idx 1 = subtitle (only if type is SUBTITLE)
                    if ph_idx == 0 or type_str == "title" or type_str == "center_title":
                        is_title = True
                    elif "subtitle" in type_str:
                        is_subtitle = True
            except (AttributeError, KeyError):
                pass

            # Fallback: if shape is at the top of the slide and large, treat as title
            if not is_title and not is_subtitle:
                try:
                    if shape.top is not None and shape.top < Emu(914400 * 1.5):
                        # Top 1.5 inches
                        if shape.height is not None and shape.height > Emu(914400 * 0.5):
                            is_title = True
                except (TypeError, AttributeError):
                    pass

            if is_title:
                slide_data["title"] = text
            elif is_subtitle:
                slide_data["subtitle"] = text
            else:
                # Body text: collect paragraphs
                for para in tf.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        level = 0
                        try:
                            level = para.level
                        except (AttributeError, KeyError):
                            pass
                        slide_data["body"].append(para_text)
                        slide_data["paragraphs"].append({
                            "text": para_text,
                            "level": level,
                        })

        # Extract images
        if "PICTURE" in shape_type or "IMAGE" in shape_type:
            image_info = {
                "path": "",
                "left": emu_to_px(shape.left),
                "top": emu_to_px(shape.top),
                "width": emu_to_px(shape.width),
                "height": emu_to_px(shape.height),
            }
            try:
                image_info["path"] = shape.image.filename or ""
            except (AttributeError, KeyError):
                try:
                    image_info["path"] = shape.image.blob.filename or "image"
                except (AttributeError, KeyError):
                    image_info["path"] = "image_{}".format(len(slide_data["images"]))
            slide_data["images"].append(image_info)

    # Extract speaker notes
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text
    except (AttributeError, KeyError):
        pass

    return slide_data


def import_pptx(pptx_path, output_path):
    """Main import function: read pptx, extract data, write JSON."""
    if not os.path.exists(pptx_path):
        print("ERROR: File not found: {}".format(pptx_path), file=sys.stderr)
        sys.exit(1)

    print("Importing: {}".format(pptx_path))

    # Extract theme colors and fonts from the raw zip
    theme = extract_theme_from_pptx(pptx_path)
    print("  Theme colors found: {}".format(len(theme["colors"])))
    print("  Theme fonts found: {}".format(len(theme["fonts"])))

    # Open with python-pptx and extract slide data
    prs = Presentation(pptx_path)
    print("  Slides found: {}".format(len(prs.slides)))

    slides_data = []
    for index, slide in enumerate(prs.slides):
        slide_data = extract_slide_data(slide, index)
        slides_data.append(slide_data)
        title_preview = slide_data["title"][:50] if slide_data["title"] else "(no title)"
        print("    Slide {}: {} [{}]".format(index, title_preview, slide_data["layout"]))

    result = {
        "source_file": os.path.basename(pptx_path),
        "slide_count": len(slides_data),
        "theme": theme,
        "slides": slides_data,
    }

    # Ensure output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print("  Output written: {}".format(output_path))
    print("  Total slides: {}".format(len(slides_data)))
    print("  Total images: {}".format(sum(len(s["images"]) for s in slides_data)))
    print("  Total notes: {}".format(sum(1 for s in slides_data if s["notes"])))


def main():
    parser = argparse.ArgumentParser(
        description="Import a .pptx file and convert to DeckMason JSON format."
    )
    parser.add_argument("input", help="Path to the .pptx file to import")
    parser.add_argument(
        "--output", "-o",
        default="decks/input_deck.json",
        help="Output JSON file path (default: decks/input_deck.json)"
    )
    args = parser.parse_args()

    import_pptx(args.input, args.output)


if __name__ == "__main__":
    main()
